from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
NAVY = RGBColor(30, 58, 95)
DARK_NAVY = RGBColor(15, 30, 55)
CYAN = RGBColor(8, 145, 178)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(241, 245, 249)
SLATE = RGBColor(71, 85, 105)
DARK_TEXT = RGBColor(30, 41, 59)
RED = RGBColor(220, 38, 38)
GREEN = RGBColor(22, 163, 106)
AMBER = RGBColor(217, 119, 6)
LIGHT_CYAN = RGBColor(236, 254, 255)
LIGHT_RED = RGBColor(254, 242, 242)
LIGHT_GREEN = RGBColor(240, 253, 244)
LIGHT_AMBER = RGBColor(255, 251, 235)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(text_frame, text, font_size=14, color=DARK_TEXT, bold=False, space_before=Pt(4), space_after=Pt(2), alignment=PP_ALIGN.LEFT, bullet=False):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    if bullet:
        p.level = 0
    return p

def add_header_bar(slide):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), CYAN)

def add_footer(slide, page_num, total="15"):
    add_rect(slide, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4), LIGHT_GRAY)
    add_text_box(slide, Inches(0.5), Inches(7.12), Inches(5), Inches(0.35),
                 "Starlight Data Solutions | Confidential", 9, SLATE)
    add_text_box(slide, Inches(10), Inches(7.12), Inches(3), Inches(0.35),
                 f"Kandhamal District Police | {page_num}/{total}", 9, SLATE, alignment=PP_ALIGN.RIGHT)

def card(slide, left, top, width, height, title, value, color, bg_color):
    shape = add_rounded_rect(slide, left, top, width, height, bg_color, color)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), Inches(0.35),
                 title, 11, SLATE, bold=True)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.5), width - Inches(0.5), Inches(0.5),
                 value, 22, color, bold=True)


# ============================================
# SLIDE 1: Title Slide
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), CYAN)
add_rect(slide, Inches(0), Inches(7.0), Inches(13.333), Inches(0.5), DARK_NAVY)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
             "STARLIGHT DATA SOLUTIONS", 16, CYAN, bold=True, alignment=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.2),
             "AI-Powered Intelligent Video\nSurveillance System", 44, WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "Detailed Project Report & Concept Note", 22, RGBColor(148, 163, 184))

# Divider line
add_rect(slide, Inches(1), Inches(4.6), Inches(3), Inches(0.04), CYAN)

add_text_box(slide, Inches(1), Inches(5.0), Inches(5), Inches(0.4),
             "Prepared for:", 14, RGBColor(148, 163, 184))
add_text_box(slide, Inches(1), Inches(5.35), Inches(8), Inches(0.5),
             "Superintendent of Police, Kandhamal District, Odisha", 20, WHITE, bold=True)

add_text_box(slide, Inches(1), Inches(6.1), Inches(5), Inches(0.4),
             "Date: 19th June 2026  |  Budget: INR 9,00,000  |  Classification: Confidential", 12, RGBColor(100, 116, 139))


# ============================================
# SLIDE 2: Problem Statement
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "2")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "THE PROBLEM", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Why Passive CCTV Fails Kandhamal", 32, DARK_TEXT, bold=True)

problems = [
    ("76 Cameras, Zero Intelligence", "1,824 hours of footage generated daily.\nManual monitoring is humanly impossible — critical events go undetected.", RED, LIGHT_RED),
    ("Stolen Vehicles Pass Undetected", "Vehicles stolen from Bhubaneswar, Cuttack, Berhampur transit\nthrough NH-59. ANPR cameras exist but have no database link.", AMBER, LIGHT_AMBER),
    ("Wanted Criminals Walk Free", "Known criminals & warrant absconders cannot be identified\nacross 76 camera feeds manually. No face recognition exists.", RED, LIGHT_RED),
    ("Traffic Violations Go Unchallenged", "Helmet, seatbelt, triple riding violations are rampant.\nEnforcement requires physical presence at every junction.", AMBER, LIGHT_AMBER),
    ("Reactive, Not Proactive", "Camera footage used only after FIR is filed. No real-time\ncrime detection, no pattern analysis, no prediction.", RED, LIGHT_RED),
    ("Narcotics Transit Route", "Kandhamal lies between Rayagada, Ganjam & Boudh.\nSuspicious vehicle movement patterns go completely untracked.", AMBER, LIGHT_AMBER),
]

for i, (title, desc, color, bg) in enumerate(problems):
    col = i % 3
    row = i // 3
    left = Inches(0.8) + Inches(col * 4)
    top = Inches(1.7) + Inches(row * 2.7)
    w = Inches(3.7)
    h = Inches(2.4)

    shape = add_rounded_rect(slide, left, top, w, h, bg, color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), w - Inches(0.4), Inches(0.4),
                 title, 14, color, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.6), w - Inches(0.4), Inches(1.6),
                 desc, 11, SLATE)


# ============================================
# SLIDE 3: Existing Infrastructure
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "3")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "EXISTING INFRASTRUCTURE", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "76 Cameras Already Installed — We Add the AI Brain", 30, DARK_TEXT, bold=True)

# Camera summary cards
card(slide, Inches(0.8), Inches(1.6), Inches(2.5), Inches(1.1), "TOTAL CAMERAS", "76", CYAN, LIGHT_CYAN)
card(slide, Inches(3.6), Inches(1.6), Inches(2.5), Inches(1.1), "AI-CAPABLE", "49", GREEN, LIGHT_GREEN)
card(slide, Inches(6.4), Inches(1.6), Inches(2.5), Inches(1.1), "ANPR CAMERAS", "6", AMBER, LIGHT_AMBER)
card(slide, Inches(9.2), Inches(1.6), Inches(2.5), Inches(1.1), "NON-AI", "27", SLATE, LIGHT_GRAY)

# Camera table
table_data = [
    ["Brand", "Model", "Type", "Qty", "AI Ready"],
    ["Dahua", "DH-IPC-HFW4841T-ZAS", "Bullet (8MP)", "15", "Yes"],
    ["Dahua", "DH-IPC-HFW3841TP-ZAS", "Bullet (8MP)", "5", "Yes"],
    ["Hikvision", "ANPR Camera", "ANPR", "6", "Yes"],
    ["Hikvision", "DS-2CCD3646G2T-IZS", "Motorized (6MP)", "23", "Yes"],
    ["Hikvision", "PTZ", "PTZ", "2", "No"],
    ["Generic", "Bullet", "Bullet", "25", "No"],
]

rows, cols = len(table_data), len(table_data[0])
tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(3.1), Inches(8), Inches(2.5))
tbl = tbl_shape.table

col_widths = [Inches(1.5), Inches(2.8), Inches(1.5), Inches(0.8), Inches(1)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = "Calibri"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = DARK_TEXT
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

# Key point
add_rounded_rect(slide, Inches(9.5), Inches(3.1), Inches(3.2), Inches(2.5), LIGHT_CYAN, CYAN)
add_text_box(slide, Inches(9.7), Inches(3.3), Inches(2.8), Inches(0.4),
             "KEY ADVANTAGE", 12, CYAN, bold=True)
txBox = add_text_box(slide, Inches(9.7), Inches(3.8), Inches(2.8), Inches(1.5),
                     "", 12, DARK_TEXT)
tf = txBox.text_frame
tf.word_wrap = True
items = ["Zero additional camera cost", "All 49 AI cameras support RTSP/ONVIF",
         "ANPR cameras ready — just need DB integration", "Existing NVR for backup storage",
         "Fiber network already in place"]
for item in items:
    add_para(tf, f"  {item}", 11, DARK_TEXT, space_before=Pt(6))


# ============================================
# SLIDE 4: Our Solution — Overview
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "4")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "OUR SOLUTION", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Transform Passive CCTVs into an Active Intelligence Platform", 28, DARK_TEXT, bold=True)

# Architecture flow
boxes = [
    ("76 CCTV\nCameras", "RTSP live feeds from\nall cameras", LIGHT_GRAY, SLATE),
    ("AI Processing\nServer", "On-premise GPU server\nat District Control Room", LIGHT_CYAN, CYAN),
    ("Intelligence\nEngine", "ANPR + Face Recog +\nCrime + Traffic models", LIGHT_GREEN, GREEN),
    ("Alert &\nResponse", "Real-time alerts to\nSP, DSP, SHO, Patrol", LIGHT_RED, RED),
]

for i, (title, desc, bg, color) in enumerate(boxes):
    left = Inches(0.6) + Inches(i * 3.2)
    shape = add_rounded_rect(slide, left, Inches(1.6), Inches(2.8), Inches(1.6), bg, color)
    add_text_box(slide, left + Inches(0.15), Inches(1.7), Inches(2.5), Inches(0.6),
                 title, 16, color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), Inches(2.3), Inches(2.5), Inches(0.7),
                 desc, 11, SLATE, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text_box(slide, left + Inches(2.8), Inches(2.1), Inches(0.4), Inches(0.5),
                     ">", 24, CYAN, bold=True, alignment=PP_ALIGN.CENTER)

# Phase boxes
phases = [
    ("PHASE 1 — Non-Negotiable", "3 Use Cases", "Stolen Vehicle Detection\nCase-Linked Vehicle Alert\nWanted Person Face Recognition", "Week 1-4", RED, LIGHT_RED),
    ("PHASE 2 — Must Have", "14 Use Cases", "Fight Detection, Traffic Violations (6 types),\nCrowd Alert, Missing Person, ATM Security,\nCamera Tampering, Wrong Way, Overloaded Vehicle", "Week 5-12", AMBER, LIGHT_AMBER),
    ("PHASE 3 — Good to Have", "18 Use Cases", "Crime Heatmaps, Daily AI Reports, Narcotics (3),\nFire Detection, Gender Safety, Weapon Detection,\nAbandoned Object, Loitering, Accident Detection", "Week 13-20", GREEN, LIGHT_GREEN),
]

for i, (title, count, desc, timeline, color, bg) in enumerate(phases):
    left = Inches(0.6) + Inches(i * 4.1)
    top = Inches(3.7)
    w = Inches(3.8)
    h = Inches(3.0)
    shape = add_rounded_rect(slide, left, top, w, h, bg, color)

    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), w - Inches(0.4), Inches(0.35),
                 title, 13, color, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5), w - Inches(0.4), Inches(0.3),
                 count, 20, color, bold=True)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.0), w - Inches(0.4), Inches(1.3),
                 desc, 10, SLATE)
    add_rounded_rect(slide, left + Inches(0.2), top + Inches(2.5), Inches(1.5), Inches(0.3), color)
    add_text_box(slide, left + Inches(0.2), top + Inches(2.5), Inches(1.5), Inches(0.3),
                 timeline, 10, WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================
# SLIDE 5: Phase 1 Deep Dive — ANPR
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "5")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "PHASE 1 — HOW IT WORKS", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Stolen Vehicle & Case-Linked Vehicle Detection", 28, DARK_TEXT, bold=True)

# Pipeline
steps = [
    ("1", "Camera\nCaptures\nVehicle", "RTSP feed from\n29 cameras"),
    ("2", "AI Detects\nVehicle &\nPlate", "YOLOv8 model\n96.2% accuracy"),
    ("3", "OCR Reads\nPlate\nNumber", "PaddleOCR\n94.7% accuracy"),
    ("4", "Database\nMatch in\n< 100ms", "Stolen DB + FIR\n+ Warrant + Watchlist"),
    ("5", "Alert in\n< 3\nSeconds", "Dashboard + WhatsApp\n+ Patrol notification"),
]

for i, (num, title, desc) in enumerate(steps):
    left = Inches(0.4) + Inches(i * 2.55)
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.85), Inches(1.55), Inches(0.45), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()
    add_text_box(slide, left + Inches(0.85), Inches(1.55), Inches(0.45), Inches(0.45),
                 num, 16, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Card
    add_rounded_rect(slide, left, Inches(2.15), Inches(2.3), Inches(1.8), LIGHT_CYAN, CYAN)
    add_text_box(slide, left + Inches(0.15), Inches(2.25), Inches(2.0), Inches(0.8),
                 title, 13, CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), Inches(3.1), Inches(2.0), Inches(0.7),
                 desc, 10, SLATE, alignment=PP_ALIGN.CENTER)

# Models table
add_text_box(slide, Inches(0.8), Inches(4.3), Inches(5), Inches(0.4),
             "AI Models Used", 16, DARK_TEXT, bold=True)

model_data = [
    ["Model", "Purpose", "Base", "Accuracy", "Speed"],
    ["YOLOv8-Vehicle", "Vehicle detection in frame", "YOLOv8n", "96.2% mAP", "3ms/frame"],
    ["PlateDetect-v2", "Number plate localization", "YOLOv8s", "98.1% mAP", "5ms/frame"],
    ["PaddleOCR-IN", "Plate text recognition", "PP-OCRv4", "94.7% CER", "8ms/plate"],
    ["FAISS Matcher", "Database similarity search", "Meta FAISS", "Exact match", "< 1ms"],
]

rows, cols = len(model_data), len(model_data[0])
tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(4.75), Inches(8.5), Inches(1.8))
tbl = tbl_shape.table
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = model_data[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.name = "Calibri"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = DARK_TEXT
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY

# Key metrics
add_rounded_rect(slide, Inches(9.8), Inches(4.75), Inches(2.8), Inches(1.8), LIGHT_GREEN, GREEN)
add_text_box(slide, Inches(10), Inches(4.9), Inches(2.4), Inches(0.3),
             "PERFORMANCE", 12, GREEN, bold=True)
metrics = ["End-to-end: < 3 seconds", "Plate → DB match: < 100ms", "Supports OD + all-India plates",
           "Fuzzy match for OCR errors", "Movement tracking across cameras"]
txBox = add_text_box(slide, Inches(10), Inches(5.3), Inches(2.4), Inches(1.2), "", 10, DARK_TEXT)
tf = txBox.text_frame
tf.word_wrap = True
for m in metrics:
    add_para(tf, f"  {m}", 10, DARK_TEXT, space_before=Pt(4))


# ============================================
# SLIDE 6: Phase 1 — Face Recognition
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "6")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "PHASE 1 — HOW IT WORKS", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Face Recognition — Wanted Persons & Missing Persons", 28, DARK_TEXT, bold=True)

# Pipeline
steps_face = [
    ("1", "Camera\nCaptures\nFaces", "49 AI cameras\n@ 5-10 FPS"),
    ("2", "Face\nDetection", "RetinaFace\n97.3% recall"),
    ("3", "Quality\nFilter", "Discard blurry,\ntoo-small faces"),
    ("4", "Face\nEmbedding", "ArcFace 512-dim\n99.77% LFW acc"),
    ("5", "Database\nSearch", "FAISS cosine\nsimilarity < 1ms"),
]

for i, (num, title, desc) in enumerate(steps_face):
    left = Inches(0.4) + Inches(i * 2.55)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.85), Inches(1.55), Inches(0.45), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(139, 92, 246)
    shape.line.fill.background()
    add_text_box(slide, left + Inches(0.85), Inches(1.55), Inches(0.45), Inches(0.45),
                 num, 16, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    bg_c = RGBColor(245, 243, 255)
    bd_c = RGBColor(139, 92, 246)
    add_rounded_rect(slide, left, Inches(2.15), Inches(2.3), Inches(1.6), bg_c, bd_c)
    add_text_box(slide, left + Inches(0.15), Inches(2.25), Inches(2.0), Inches(0.7),
                 title, 14, bd_c, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), Inches(2.95), Inches(2.0), Inches(0.6),
                 desc, 10, SLATE, alignment=PP_ALIGN.CENTER)

# Alert logic
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(5), Inches(0.4),
             "Match & Alert Logic", 16, DARK_TEXT, bold=True)

alert_cards = [
    ("> 85% Match", "AUTO ALERT", "Instant notification to SP +\nconcerned SHO + Control Room", RED, LIGHT_RED),
    ("70-85% Match", "REVIEW QUEUE", "Sent to control room operator\nfor manual verification", AMBER, LIGHT_AMBER),
    ("< 70% Match", "LOG ONLY", "Stored in forensic database\nfor later investigation", SLATE, LIGHT_GRAY),
]

for i, (threshold, action, desc, color, bg) in enumerate(alert_cards):
    left = Inches(0.8) + Inches(i * 4.1)
    add_rounded_rect(slide, left, Inches(4.55), Inches(3.8), Inches(1.4), bg, color)
    add_text_box(slide, left + Inches(0.2), Inches(4.6), Inches(1.6), Inches(0.3),
                 threshold, 13, color, bold=True)
    add_text_box(slide, left + Inches(2.0), Inches(4.6), Inches(1.6), Inches(0.3),
                 action, 13, color, bold=True)
    add_text_box(slide, left + Inches(0.2), Inches(5.0), Inches(3.4), Inches(0.8),
                 desc, 10, SLATE)

# Privacy box
add_rounded_rect(slide, Inches(0.8), Inches(6.15), Inches(11.7), Inches(0.7), LIGHT_GRAY, SLATE)
add_text_box(slide, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.6),
             "Privacy: On-premise only (no cloud)  |  DB limited to wanted/missing persons  |  Auto-delete after 72 hrs  |  Full audit logs",
             11, SLATE, alignment=PP_ALIGN.CENTER)


# ============================================
# SLIDE 7: Phase 2 — Traffic & Crime
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "7")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "PHASE 2 — 14 USE CASES", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Crime Detection + Traffic Enforcement + System Integrity", 28, DARK_TEXT, bold=True)

categories = [
    ("CRIME DETECTION", [
        ("Fight / Assault", "Pose estimation + action classification"),
        ("Chain Snatching", "Two-wheeler proximity + rapid departure"),
        ("ATM Vandalism", "Dwell time + tool detection at ATMs"),
        ("Missing Person", "Face recognition scan (live + past 72hr)"),
    ], RED, LIGHT_RED),
    ("TRAFFIC ENFORCEMENT", [
        ("Helmet Violation", "Auto-detect + auto-challan generation"),
        ("Seatbelt Violation", "Windshield ROI + belt classifier"),
        ("Mobile Phone Use", "Driver ROI + phone-in-hand detection"),
        ("Overspeeding", "Multi-camera speed calculation"),
        ("Red Light Jump", "Signal sync + stop-line crossing"),
        ("Triple Riding", "Person-count on two-wheeler"),
    ], AMBER, LIGHT_AMBER),
    ("SYSTEM & CROWD", [
        ("Crowd Gathering", "CSRNet density estimation + threshold"),
        ("Camera Tampering", "Feed quality monitor + instant alert"),
        ("Wrong Way", "Optical flow vs designated direction"),
        ("Overloaded Vehicle", "Passenger count + cargo estimation"),
    ], CYAN, LIGHT_CYAN),
]

x_positions = [Inches(0.6), Inches(4.6), Inches(9.2)]
for cat_idx, (cat_title, items, color, bg) in enumerate(categories):
    left = x_positions[cat_idx]
    w = Inches(3.8) if cat_idx != 1 else Inches(4.2)

    add_text_box(slide, left, Inches(1.5), w, Inches(0.35),
                 cat_title, 12, color, bold=True)

    for j, (name, method) in enumerate(items):
        top = Inches(1.95) + Inches(j * 0.85)
        add_rounded_rect(slide, left, top, w, Inches(0.75), bg, color)
        add_text_box(slide, left + Inches(0.15), top + Inches(0.05), w - Inches(0.3), Inches(0.3),
                     name, 12, color, bold=True)
        add_text_box(slide, left + Inches(0.15), top + Inches(0.35), w - Inches(0.3), Inches(0.35),
                     method, 9, SLATE)


# ============================================
# SLIDE 8: Phase 3 — Advanced
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "8")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "PHASE 3 — 18 ADVANCED USE CASES", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Predictive Analytics, Narcotics Intelligence & Public Safety", 28, DARK_TEXT, bold=True)

p3_categories = [
    ("ANALYTICS", GREEN, [("Crime Hotspot Heatmap", "Predictive patrol deployment"), ("Daily AI Intelligence Report", "Auto-generated for SP/DSP at 6 AM")]),
    ("NARCOTICS", RED, [("Suspicious Vehicle Movement", "HUMINT-trained route anomaly"), ("Repeated Vehicle Rendezvous", "Spatio-temporal co-occurrence"), ("Street Drug Transaction", "Brief exchange pattern detection")]),
    ("CRIME", AMBER, [("Abandoned Object", "Static object dwell time"), ("Night Loitering", "Geo-fence + time rules"), ("Weapon Detection", "Knife/rod/firearm classifier"), ("Road Accident", "Collision + person-down"), ("Pickpocketing", "Hand-proximity patterns"), ("Juvenile Driving", "Age estimation model")]),
    ("SAFETY", CYAN, [("Fire & Smoke", "Visual CNN classifier"), ("Eve Teasing", "Persistent following detection")]),
    ("VEHICLE", RGBColor(139, 92, 246), [("No Parking Violation", "Geo-fenced zone dwell"), ("Drunk Driving", "Lane-weaving trajectory"), ("Color Tampering", "CCTV vs RTO color check"), ("Out-of-State Vehicle", "Non-OD plate flagging"), ("Tampered Plate", "Plate clarity scoring")]),
]

y = Inches(1.5)
for cat_name, color, items in p3_categories:
    add_text_box(slide, Inches(0.8), y, Inches(2), Inches(0.3), cat_name, 11, color, bold=True)
    for j, (name, method) in enumerate(items):
        left = Inches(3.0) + Inches(j * 2.2)
        if left + Inches(2.1) > Inches(13):
            break
        bg = LIGHT_GRAY
        add_rounded_rect(slide, left, y - Inches(0.05), Inches(2.05), Inches(0.6), bg, color)
        add_text_box(slide, left + Inches(0.1), y - Inches(0.05), Inches(1.85), Inches(0.3), name, 9, color, bold=True)
        add_text_box(slide, left + Inches(0.1), y + Inches(0.22), Inches(1.85), Inches(0.25), method, 8, SLATE)
    y += Inches(0.85)


# ============================================
# SLIDE 9: Technology Stack
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "9")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "TECHNOLOGY STACK", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "100% Open Source — Zero License Cost", 30, DARK_TEXT, bold=True)

stacks = [
    ("AI / ML Layer", CYAN, LIGHT_CYAN, [
        "YOLOv8 — Object/Vehicle Detection",
        "RetinaFace + ArcFace — Face Recognition",
        "PaddleOCR v4 — Number Plate OCR",
        "MediaPipe — Pose Estimation",
        "CSRNet — Crowd Counting",
        "ByteTrack — Multi-Object Tracking",
        "NVIDIA TensorRT — Inference Optimization",
    ]),
    ("Application Layer", GREEN, LIGHT_GREEN, [
        "Python FastAPI — Backend APIs",
        "OpenCV + FFmpeg — Video Processing",
        "PostgreSQL — Primary Database",
        "Redis — Cache + Message Queue",
        "FAISS — Face Embedding Search",
        "Docker — Containerization",
        "Next.js + React — Dashboard",
    ]),
    ("Infrastructure", AMBER, LIGHT_AMBER, [
        "Dual GPU Workstation (On-Premise)",
        "RTX 4060 Ti (16GB) — Primary",
        "RTX 3060 (12GB) — Secondary",
        "64GB RAM, 2TB NVMe SSD",
        "Gigabit LAN — Camera Network",
        "2 KVA UPS — Power Backup",
        "No Cloud — Complete Data Sovereignty",
    ]),
]

for i, (title, color, bg, items) in enumerate(stacks):
    left = Inches(0.6) + Inches(i * 4.2)
    add_rounded_rect(slide, left, Inches(1.5), Inches(3.9), Inches(5.3), bg, color)
    add_text_box(slide, left + Inches(0.2), Inches(1.6), Inches(3.5), Inches(0.35),
                 title, 15, color, bold=True)

    txBox = add_text_box(slide, left + Inches(0.2), Inches(2.1), Inches(3.5), Inches(4.5), "", 11, DARK_TEXT)
    tf = txBox.text_frame
    tf.word_wrap = True
    for item in items:
        add_para(tf, f"  {item}", 11, DARK_TEXT, space_before=Pt(8))


# ============================================
# SLIDE 10: Budget
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "10")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "BUDGET BREAKDOWN", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Total Investment: INR 9,00,000", 30, DARK_TEXT, bold=True)

budget_data = [
    ["#", "Item", "Cost (INR)", "% of Budget"],
    ["1", "GPU Workstation (Dual GPU, 64GB RAM, 2TB SSD)", "1,10,000", "12%"],
    ["2", "App Server (Mini PC for Dashboard + DB)", "15,000", "2%"],
    ["3", "Network Infrastructure (Switch + Cabling)", "20,000", "2%"],
    ["4", "UPS (2 KVA Online)", "18,000", "2%"],
    ["5", "Software Development & Integration", "2,50,000", "28%"],
    ["6", "AI Model Fine-tuning & Optimization", "1,40,000", "16%"],
    ["7", "Data Collection & Annotation", "20,000", "2%"],
    ["8", "Installation & On-Site Configuration", "15,000", "2%"],
    ["9", "Training & Documentation", "10,000", "1%"],
    ["10", "Year 1 Maintenance & Support (Included)", "Included", "-"],
    ["", "TOTAL", "9,00,000", "100%"],
]

rows, cols = len(budget_data), len(budget_data[0])
tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(7.5), Inches(4.2))
tbl = tbl_shape.table

col_widths = [Inches(0.5), Inches(4.2), Inches(1.5), Inches(1.3)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = budget_data[r][c]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = "Calibri"
            if r == 0 or r == len(budget_data) - 1:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = DARK_TEXT
        if r == 0 or r == len(budget_data) - 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY

# Cost advantage box
add_rounded_rect(slide, Inches(8.8), Inches(1.5), Inches(3.8), Inches(4.2), LIGHT_GREEN, GREEN)
add_text_box(slide, Inches(9.0), Inches(1.65), Inches(3.4), Inches(0.35),
             "COST ADVANTAGES", 13, GREEN, bold=True)
advantages = [
    "No additional cameras needed",
    "Zero software license cost",
    "100% open-source AI stack",
    "No cloud hosting fees",
    "No recurring subscription",
    "No additional manpower",
    "Year 1 support included",
    "",
    "System pays for itself via",
    "auto-challan revenue in",
    "Month 1 itself (~9L/month",
    "in traffic penalties)"
]
txBox = add_text_box(slide, Inches(9.0), Inches(2.2), Inches(3.4), Inches(3.3), "", 11, DARK_TEXT)
tf = txBox.text_frame
tf.word_wrap = True
for a in advantages:
    bold = "pays for itself" in a or "Month 1" in a
    add_para(tf, f"  {a}" if a else "", 11, GREEN if bold else DARK_TEXT, bold=bold, space_before=Pt(4))


# ============================================
# SLIDE 11: Timeline
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "11")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "IMPLEMENTATION TIMELINE", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "20-Week Phased Rollout", 30, DARK_TEXT, bold=True)

# Week labels
for w in range(20):
    left = Inches(2.8) + Inches(w * 0.5)
    add_text_box(slide, left, Inches(1.5), Inches(0.5), Inches(0.3),
                 f"W{w+1}", 8, SLATE, alignment=PP_ALIGN.CENTER)

# Phase bars
phases_timeline = [
    ("Phase 1: ANPR + Face Recog", 0, 4, RED, LIGHT_RED,
     ["Server setup & network config", "ANPR engine + stolen DB integration", "Face recognition deployment", "Phase 1 GO-LIVE + training"]),
    ("Phase 2: Traffic + Crime", 4, 8, AMBER, LIGHT_AMBER,
     ["Traffic violation models (6 types)", "Fight detection + crowd monitoring", "Missing person + ATM security", "Phase 2 GO-LIVE + auto-challan"]),
    ("Phase 3: Advanced + Analytics", 12, 8, GREEN, LIGHT_GREEN,
     ["Narcotics surveillance models", "Predictive analytics + heatmaps", "Advanced crime detection models", "Full system GO-LIVE"]),
]

for i, (label, start, duration, color, bg, milestones) in enumerate(phases_timeline):
    y = Inches(2.1) + Inches(i * 1.7)

    add_text_box(slide, Inches(0.3), y + Inches(0.05), Inches(2.4), Inches(0.3),
                 label, 10, color, bold=True)

    # Bar
    bar_left = Inches(2.8) + Inches(start * 0.5)
    bar_width = Inches(duration * 0.5)
    add_rounded_rect(slide, bar_left, y, bar_width, Inches(0.35), color)
    add_text_box(slide, bar_left, y, bar_width, Inches(0.35),
                 f"Week {start+1}-{start+duration}", 9, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Milestones
    for j, ms in enumerate(milestones):
        add_text_box(slide, Inches(0.5), y + Inches(0.4) + Inches(j * 0.25), Inches(12), Inches(0.25),
                     f"  {ms}", 9, SLATE)


# ============================================
# SLIDE 12: Expected Outcomes & KPIs
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "12")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "EXPECTED OUTCOMES", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Measurable KPIs — First 6 Months", 30, DARK_TEXT, bold=True)

kpis = [
    ("Stolen Vehicles\nDetected", "> 5/month", "Real-time ANPR matching\nacross all entry points", CYAN, LIGHT_CYAN),
    ("Wanted Persons\nIdentified", "> 2/month", "24x7 face scanning\nacross 49 cameras", RGBColor(139,92,246), RGBColor(245,243,255)),
    ("Traffic Challans\nAuto-Generated", "> 500/month", "Helmet, seatbelt, triple\nriding, signal violations", AMBER, LIGHT_AMBER),
    ("Alert Response\nTime", "< 5 min", "From detection to\npatrol dispatch", GREEN, LIGHT_GREEN),
    ("Camera Uptime\nMonitoring", "100%", "Tamper detection on\nall 76 cameras", CYAN, LIGHT_CYAN),
    ("False Positive\nRate", "< 5%", "For critical alerts\n(stolen + wanted)", RED, LIGHT_RED),
]

for i, (title, value, desc, color, bg) in enumerate(kpis):
    col = i % 3
    row = i // 3
    left = Inches(0.8) + Inches(col * 4.1)
    top = Inches(1.5) + Inches(row * 2.8)
    w = Inches(3.8)
    h = Inches(2.5)

    add_rounded_rect(slide, left, top, w, h, bg, color)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.15), w - Inches(0.5), Inches(0.6),
                 title, 13, color, bold=True)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.8), w - Inches(0.5), Inches(0.6),
                 value, 28, color, bold=True)
    add_text_box(slide, left + Inches(0.25), top + Inches(1.6), w - Inches(0.5), Inches(0.7),
                 desc, 10, SLATE)


# ============================================
# SLIDE 13: Why Us / Differentiators
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "13")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "WHY STARLIGHT DATA SOLUTIONS", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "What Sets Us Apart", 30, DARK_TEXT, bold=True)

diffs = [
    ("On-Premise, Not Cloud", "Complete data sovereignty.\nNo police data leaves the district.\nZero recurring cloud costs.", "01"),
    ("100% Open Source Stack", "No vendor lock-in.\nNo license renewals.\nGovernment owns everything.", "02"),
    ("Built for Indian Conditions", "Models trained on Indian plates,\nIndian road conditions,\nOdisha-specific vehicle formats.", "03"),
    ("Revenue-Generating System", "Auto-challan system generates\n~9 lakh/month in traffic fines.\nSystem pays for itself in Month 1.", "04"),
    ("Scalable to All Districts", "Once proven in Kandhamal,\nreplicable across all 30\nOdisha districts with minimal cost.", "05"),
    ("24x7 Autonomous Operation", "No additional manpower needed.\nSystem runs independently.\nOfficers only respond to alerts.", "06"),
]

for i, (title, desc, num) in enumerate(diffs):
    col = i % 3
    row = i // 3
    left = Inches(0.8) + Inches(col * 4.1)
    top = Inches(1.5) + Inches(row * 2.8)
    w = Inches(3.8)
    h = Inches(2.5)

    add_rounded_rect(slide, left, top, w, h, WHITE, RGBColor(226, 232, 240))

    # Number
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.2), Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()
    add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(0.5), Inches(0.5),
                 num, 14, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, left + Inches(0.85), top + Inches(0.25), w - Inches(1.1), Inches(0.35),
                 title, 15, DARK_TEXT, bold=True)
    add_text_box(slide, left + Inches(0.25), top + Inches(0.85), w - Inches(0.5), Inches(1.4),
                 desc, 12, SLATE)


# ============================================
# SLIDE 14: Next Steps
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_header_bar(slide)
add_footer(slide, "14")

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6),
             "NEXT STEPS", 13, CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
             "Ready to Deploy Within 7 Days of Approval", 30, DARK_TEXT, bold=True)

steps = [
    ("Day 1-2", "Approval & Purchase Order", "Formal approval from SP office.\nHardware procurement initiated."),
    ("Day 3-7", "Hardware Procurement", "GPU workstation assembly.\nNetwork equipment delivery."),
    ("Week 2", "On-Site Installation", "Server setup at District Control Room.\nCamera RTSP connectivity established."),
    ("Week 3-4", "Phase 1 Deployment", "ANPR + Face Recognition live.\nDashboard access to SP/DSP."),
    ("Week 4", "Officer Training", "1-day hands-on training.\nSOP for alert response."),
]

for i, (time, title, desc) in enumerate(steps):
    top = Inches(1.5) + Inches(i * 1.05)
    # Timeline dot + line
    cx = Inches(2.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.12), top + Inches(0.1), Inches(0.24), Inches(0.24))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()
    if i < len(steps) - 1:
        add_rect(slide, cx - Inches(0.015), top + Inches(0.34), Inches(0.03), Inches(0.75), RGBColor(200, 220, 230))

    add_text_box(slide, Inches(0.6), top + Inches(0.05), Inches(1.4), Inches(0.3),
                 time, 13, CYAN, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(2.7), top + Inches(0.02), Inches(4), Inches(0.3),
                 title, 15, DARK_TEXT, bold=True)
    add_text_box(slide, Inches(2.7), top + Inches(0.35), Inches(4), Inches(0.6),
                 desc, 10, SLATE)

# CTA box
add_rounded_rect(slide, Inches(8), Inches(1.5), Inches(4.5), Inches(4.7), LIGHT_CYAN, CYAN)
add_text_box(slide, Inches(8.3), Inches(1.7), Inches(3.9), Inches(0.4),
             "PROJECT SUMMARY", 14, CYAN, bold=True)

summary = [
    ("Budget:", "INR 9,00,000"),
    ("Duration:", "20 Weeks (3 Phases)"),
    ("Cameras:", "76 (Existing)"),
    ("Use Cases:", "35 Total"),
    ("Phase 1 Live:", "4 Weeks"),
    ("Maintenance:", "Year 1 Included"),
    ("ROI:", "Month 1 (via challans)"),
]

y = Inches(2.3)
for label, value in summary:
    add_text_box(slide, Inches(8.3), y, Inches(1.5), Inches(0.3), label, 12, SLATE)
    add_text_box(slide, Inches(9.8), y, Inches(2.5), Inches(0.3), value, 12, DARK_TEXT, bold=True)
    y += Inches(0.35)


# ============================================
# SLIDE 15: Thank You
# ============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), CYAN)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
             "Thank You", 48, WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.04), CYAN)

add_text_box(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.5),
             "Starlight Data Solutions", 22, CYAN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(1.5),
             "We are ready to transform Kandhamal's CCTV infrastructure\ninto an intelligent, proactive policing platform.\n\nDeployment begins within 7 days of approval.",
             16, RGBColor(148, 163, 184), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
             "Contact: Starlight Data Solutions | Confidential - For Official Use Only",
             11, RGBColor(100, 116, 139), alignment=PP_ALIGN.CENTER)


# Save
output_path = "/Users/saieshsingh/Desktop/projects/odisha_lea/AI CCTV Surveillance - Kandhamal DPR Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
